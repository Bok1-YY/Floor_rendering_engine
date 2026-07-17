# -*- coding: utf-8 -*-
"""导出层 —— 记录页 HTML 对照文档与客户提案 PPTX(单记录/收藏夹)。

从 records.py 迁出;函数体逐字未动。依赖 records 的 CRUD(单向,无环)。
"""
import base64 as b64mod
import html
import io
import json
import os
import time

from PIL import Image

from .config import MAIN_OUTPUT_DIR, logger, get_pptx_branding
from .records import (
    load_records_file, save_records_file, scan_json_files, get_record_labels,
    record_file_lock, safe_output_path,
)


def _html_b64(obj, file_key, b64_key):
    """HTML 导出取图：优先从文件读回 base64(保持 HTML 自包含)，回退内联 base64。"""
    rel = obj.get(file_key, '')
    p = safe_output_path(rel)
    if p:
        try:
            with open(p, 'rb') as f:
                return b64mod.b64encode(f.read()).decode()
        except Exception as e:
            logger.warning(f"[导出] 读取图片失败 {rel}: {e}")
    return obj.get(b64_key, '')


def export_html_from_json(json_path):
    records = load_records_file(json_path)
    if not records: return "❌ 没有找到记录"
    html_path = json_path.replace('_记录.json', '_导出.html')
    entries = []
    for r in records:
        _sb = _html_b64(r, 'sample_image_file', 'sample_image_b64')
        sample_tag = (f'<img src="data:image/jpeg;base64,{_sb}" '
                      f'style="width:180px;border-radius:6px;margin-bottom:10px;display:block;" />'
                      ) if _sb else ''
        results_html = ''
        for i, res in enumerate(r.get('results', [])):
            cmt = (f'<div style="margin-top:8px;padding:8px 12px;background:#fff8f0;border-left:3px solid #e8874a;border-radius:4px;">'
                   f'<p style="margin:0;font-size:0.82em;color:#7a4010;">💬 备注：{html.escape(res.get("comment",""))}</p></div>'
                   ) if res.get('comment') else ''
            review_bits = []
            if res.get('best'):
                review_bits.append('最佳图')
            _status = {'pass': '通过', 'backup': '备选', 'rejected': '淘汰', 'unreviewed': '未评'}.get(
                res.get('review_status', 'unreviewed'), '未评')
            if _status != '未评':
                review_bits.append(_status)
            if res.get('review_tags'):
                review_bits.append('标签：' + '、'.join(html.escape(str(t)) for t in res.get('review_tags', [])))
            if res.get('review_note'):
                review_bits.append('评审备注：' + html.escape(res.get('review_note', '')))
            review = (f'<div style="margin-top:8px;padding:8px 12px;background:#f3fbf8;border-left:3px solid #2e8c7e;border-radius:4px;">'
                      f'<p style="margin:0;font-size:0.82em;color:#235f55;">{" ｜ ".join(review_bits)}</p></div>'
                      ) if review_bits else ''
            _rb = _html_b64(res, 'result_image_file', 'result_image_b64')
            results_html += (f'<div style="margin-top:12px;padding:12px;background:#fdf8f4;border-radius:8px;">'
                             f'<p style="margin:0 0 8px 0;font-size:0.8em;color:#b07040;">📸 效果图 {i+1} — {res.get("result_timestamp","")}</p>'
                             f'<img src="data:image/jpeg;base64,{_rb}" style="max-width:100%;border-radius:6px;" />'
                             f'{cmt}{review}</div>')
        entries.append(f'<div style="padding:20px;border-bottom:3px solid #e8874a;margin-bottom:10px;">'
                       f'<p style="margin:0 0 8px 0;font-size:0.85em;color:#555;"><strong>🕒 {r.get("timestamp","")}</strong>'
                       f' &nbsp;|&nbsp; {r.get("workflow_mode","")} &nbsp;|&nbsp; {r.get("room_type","")}</p>'
                       f'{sample_tag}'
                       f'<pre style="white-space:pre-wrap;font-size:0.82em;background:#f8f8f8;padding:12px;border-radius:6px;">'
                       f'{html.escape(r.get("params_summary", ""))}</pre>'
                       f'{results_html}</div>')
    full_html = ('<!DOCTYPE html><html><head><meta charset="utf-8"><title>地板效果图记录</title>'
                 '<style>body{font-family:sans-serif;max-width:960px;margin:0 auto;padding:20px;}</style>'
                 f'</head><body>{"".join(entries)}</body></html>')
    with open(html_path, 'w', encoding='utf-8') as f: f.write(full_html)
    return f"✅ 已导出：{os.path.basename(html_path)}"

def append_result_to_log(img1_path, img2_path, json_path, record_id, comment1="", comment2=""):
    if not img1_path and not img2_path: return "⚠️ 请至少上传一张效果图"
    if not json_path or not record_id: return "⚠️ 请先加载一条记录"
    try:
        with record_file_lock(json_path):
            records = load_records_file(json_path)
            for r in records:
                if r.get('id') == record_id:
                    written = []; ts = time.strftime("%Y-%m-%d %H:%M:%S")
                    for img_path, comment, label in [(img1_path, comment1, "Banana2"), (img2_path, comment2, "Pro")]:
                        if img_path:
                            rel = ''
                            try:
                                rel = _rel_result_path(save_api_result_jpg(Image.open(img_path), label, json_path.replace('_记录.json', '_优化图.png')))
                            except Exception as ex:
                                logger.warning(f"[追加] 落盘失败(回退 base64): {ex}")
                            entry = {
                                'result_id': _new_result_id(),
                                'result_timestamp': ts,
                                'comment': str(comment).strip() if comment else '',
                                'model_label': label,
                            }
                            if rel:
                                entry['result_image_file'] = rel
                            else:
                                entry['result_image_b64'] = img_to_b64(img_path, max_width=1000)
                            r.setdefault('results', []).append(entry)
                            written.append(label)
                    save_records_file(json_path, records)
                    logger.info(
                        f"[记录] 手动追加效果图 json={json_path}, record={record_id}, written={written}, "
                        f"img1={img1_path}, img2={img2_path}"
                    )
                    return f"✅ 已写入：{' + '.join(written)}"
        logger.error(f"[记录] 手动追加失败，未找到记录 json={json_path}, record={record_id}")
        return "❌ 未找到对应记录"
    except Exception as e:
        logger.exception(f"[记录] 手动追加写入失败 json={json_path}, record={record_id}")
        return f"❌ 写入失败: {e}"


def collect_favorites():
    """扫描所有记录文件，汇总被收藏(favorite=True)的结果，按出图时间倒序。
    返回 [{json_path, material, record_id, record, result_id, res}, ...]。"""
    out = []
    for jp in scan_json_files():
        material = os.path.basename(jp).replace('_记录.json', '').replace('.json', '')
        try:
            recs = load_records_file(jp)
        except Exception:
            continue
        for r in recs:
            for res in r.get('results', []):
                if res.get('favorite'):
                    out.append({
                        'json_path': jp, 'material': material,
                        'record_id': r.get('id', ''), 'record': r,
                        'result_id': res.get('result_id', ''), 'res': res,
                    })
    out.sort(key=lambda x: x['res'].get('result_timestamp', ''), reverse=True)
    return out




# ── PPTX 导出（客户提案 deck）──────────────────────────────────────
def _result_image_source(res):
    """结果图供 PPTX 用的来源：优先成图文件绝对路径，回退 base64 解成 BytesIO；无图返回 None。"""
    p = safe_output_path(res.get('result_image_file', ''))
    if p:
        return p
    b64 = res.get('result_image_b64', '')
    if b64:
        try:
            return io.BytesIO(b64mod.b64decode(b64))
        except Exception as e:
            logger.warning(f"[PPTX] base64 解码失败: {e}")
    return None


def _build_pptx(items, out_path, title, branding=None):
    """items: [{'source': 路径或BytesIO, 'caption': 文本}]。16:9 PPTX：标题页 + 每项一页(图+说明)。
    python-pptx 只存字体名(微软雅黑)，PowerPoint 用系统字体渲染中文，无需嵌字体。
    branding: 可选 {'company','contact','logo_path'}（config.get_pptx_branding()）——
    标题页加 logo/公司名/联系方式，内容页右下角加公司名页脚；配置坏/图打不开一律静默降级。"""
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt, Emu
        from pptx.enum.text import PP_ALIGN
        from pptx.dml.color import RGBColor
    except Exception:
        return "❌ 未安装 python-pptx，请先运行: pip install python-pptx"
    from PIL import Image as _PILImage

    branding = branding or {}
    company = str(branding.get('company') or '').strip()
    contact = str(branding.get('contact') or '').strip()
    logo_path = str(branding.get('logo_path') or '').strip()

    prs = Presentation()
    prs.slide_width = Inches(13.333); prs.slide_height = Inches(7.5)  # 16:9
    SW, SH = int(prs.slide_width), int(prs.slide_height)
    blank = prs.slide_layouts[6]

    def _set_text(text_frame, text, size, bold, color, align=None):
        p = text_frame.paragraphs[0]
        if align is not None: p.alignment = align
        run = p.add_run(); run.text = text
        run.font.size = Pt(size); run.font.bold = bold
        run.font.name = '微软雅黑'; run.font.color.rgb = color

    # 标题页（logo 置顶居中 → 主标题 → 公司名副标题 → 联系方式底部小字）
    s0 = prs.slides.add_slide(blank)
    if logo_path and os.path.isfile(logo_path):
        try:
            with _PILImage.open(logo_path) as im:
                lw, lh = im.size
            # Pillow 给的是像素，python-pptx 使用 EMU。这里的 scale 单位是
            # EMU/像素，不能再限制为 <= 1，否则普通 1000px logo 只会得到
            # 1000 EMU（约 0.001 英寸），在封面上几乎不可见。
            max_w, max_h = int(Inches(2.4)), int(Inches(1.2))
            scale = min(max_w / max(1, lw), max_h / max(1, lh))
            pw, ph = int(lw * scale), int(lh * scale)
            s0.shapes.add_picture(logo_path, int((SW - pw) / 2), int(Inches(1.2)),
                                  width=pw, height=ph)
        except Exception as e:
            logger.warning(f"[PPTX] logo 插入失败(忽略): {e}")
    tb = s0.shapes.add_textbox(Inches(0.8), Inches(2.8), SW - Inches(1.6), Inches(1.8))
    tb.text_frame.word_wrap = True
    _set_text(tb.text_frame, title, 40, True, RGBColor(0x2b, 0x24, 0x1a), PP_ALIGN.CENTER)
    if company:
        sb = s0.shapes.add_textbox(Inches(0.8), Inches(4.4), SW - Inches(1.6), Inches(0.7))
        sb.text_frame.word_wrap = True
        _set_text(sb.text_frame, company, 18, False, RGBColor(0x6b, 0x63, 0x56), PP_ALIGN.CENTER)
    if contact:
        cb = s0.shapes.add_textbox(Inches(0.8), SH - Inches(1.0), SW - Inches(1.6), Inches(0.6))
        cb.text_frame.word_wrap = True
        _set_text(cb.text_frame, contact, 12, False, RGBColor(0x9a, 0x90, 0x82), PP_ALIGN.CENTER)

    area_w = SW - int(Inches(1.0)); area_h = SH - int(Inches(2.0))
    top0 = int(Inches(0.4))
    for it in items:
        src = it.get('source')
        if src is None:
            continue
        slide = prs.slides.add_slide(blank)
        try:
            if hasattr(src, 'seek'): src.seek(0)
            with _PILImage.open(src) as im: iw, ih = im.size
            if hasattr(src, 'seek'): src.seek(0)
        except Exception:
            iw, ih = 4, 3
        scale = min(area_w / iw, area_h / ih)
        pw = int(iw * scale); ph = int(ih * scale)
        left = int((SW - pw) / 2)
        try:
            slide.shapes.add_picture(src, left, top0, width=Emu(pw), height=Emu(ph))
        except Exception as e:
            logger.warning(f"[PPTX] 插图失败(跳过一页): {e}")
            continue
        cap = slide.shapes.add_textbox(Inches(0.5), SH - Inches(1.4), SW - Inches(1.0), Inches(1.2))
        cap.text_frame.word_wrap = True
        _set_text(cap.text_frame, it.get('caption', ''), 12, False, RGBColor(0x55, 0x55, 0x55))
        if company:
            ft = slide.shapes.add_textbox(SW - Inches(3.2), SH - Inches(0.42),
                                          Inches(3.0), Inches(0.35))
            _set_text(ft.text_frame, company, 9, False, RGBColor(0xb0, 0xa8, 0x9a), PP_ALIGN.RIGHT)

    try:
        prs.save(out_path)
    except Exception as e:
        logger.exception(f"[PPTX] 保存失败 {out_path}")
        return f"❌ 保存失败: {e}"
    logger.info(f"[PPTX] 已导出 {out_path}（{len(items)} 张）")
    return f"✅ 已导出：{os.path.basename(out_path)}"


def _result_caption(material, res, record):
    caption = " · ".join(p for p in [
        material, res.get('model_label', ''), res.get('result_timestamp', ''),
        (record.get('params_summary', '') or ''),
    ] if p)
    cmt = res.get('comment', '')
    review = []
    status = {'pass': '通过', 'backup': '备选', 'rejected': '淘汰', 'unreviewed': ''}.get(
        res.get('review_status', 'unreviewed'), '')
    if res.get('best'):
        review.append('最佳图')
    if status:
        review.append(status)
    if res.get('review_tags'):
        review.append('标签：' + '、'.join(str(t) for t in res.get('review_tags', [])))
    if res.get('review_note'):
        review.append('评审备注：' + str(res.get('review_note', '')))
    tail = (f"\n备注：{cmt}" if cmt else "")
    if review:
        tail += "\n评审：" + " ｜ ".join(review)
    return caption + tail


def export_pptx_from_json(json_path):
    """把当前材料的所有效果图导出成一份 PPTX。"""
    records = load_records_file(json_path)
    if not records:
        return "❌ 没有找到记录"
    material = os.path.basename(json_path).replace('_记录.json', '').replace('.json', '')
    items = []
    for r in records:
        for res in r.get('results', []):
            src = _result_image_source(res)
            if src is not None:
                items.append({'source': src, 'caption': _result_caption(material, res, r)})
    if not items:
        return "❌ 该文件没有可导出的效果图"
    out_path = json_path.replace('_记录.json', '_导出.pptx')
    return _build_pptx(items, out_path, f"{material} · 地板效果图", branding=get_pptx_branding())


def export_favorites_pptx():
    """把所有收藏(跨材料)的效果图合成一份客户提案 PPTX。"""
    favs = collect_favorites()
    if not favs:
        return "❌ 还没有收藏任何效果图（先在记录里点 ⭐）"
    items = []
    for f in favs:
        src = _result_image_source(f['res'])
        if src is not None:
            items.append({'source': src, 'caption': _result_caption(f['material'], f['res'], f['record'])})
    if not items:
        return "❌ 收藏的效果图都没有可用图片数据"
    out_path = os.path.join(MAIN_OUTPUT_DIR, f"收藏夹提案_{time.strftime('%Y%m%d_%H%M%S')}.pptx")
    return _build_pptx(items, out_path, "地板效果图 · 收藏提案", branding=get_pptx_branding())


