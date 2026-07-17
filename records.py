import os
import json
import time
import io
import base64 as b64mod
import hashlib
import html
import shutil
import tempfile
import threading
import dataclasses
import uuid
from typing import List, Tuple, Optional

from PIL import Image, PngImagePlugin

from .config import (
    BASE_DIR, MAIN_OUTPUT_DIR, CONFIG_FILE, UPLOAD_DIR,
    logger, _short_text, _load_config, _save_config, get_pptx_branding,
)
from .models import JobRecord, ensure_candidate_lists

# ── 记录文件并发保护 ──────────────────────────────────────────────
# 同一记录 JSON 的「读-改-写」必须串行：双模型生成时 B2/Pro 在两个 worker 线程里
# 并发 append 同一文件，无锁会丢结果。按文件路径取锁，不同素材的记录互不阻塞。
_path_locks_guard = threading.Lock()
_path_locks: dict = {}

def record_file_lock(json_path) -> threading.Lock:
    key = os.path.normcase(os.path.abspath(str(json_path)))
    with _path_locks_guard:
        lock = _path_locks.get(key)
        if lock is None:
            lock = threading.Lock()
            _path_locks[key] = lock
        return lock

def _get_json_path(image_path):
    base_name = os.path.splitext(os.path.basename(image_path))[0]
    target_dir = os.path.join(MAIN_OUTPUT_DIR, base_name)
    os.makedirs(target_dir, exist_ok=True)
    return os.path.join(target_dir, f"{base_name}_记录.json"), base_name, target_dir


# ── 渲染队列持久化：重启后恢复已完成的卡片（图片本就在盘上）──────────────
# 从 webui 下沉到此（纯文件 IO）。webui 保留薄包装 _persist_jobs() 负责加锁取 _job_history 快照后调本函数。
QUEUE_STATE_FILE = os.path.join(MAIN_OUTPUT_DIR, '.queue_state.json')
QUEUE_PERSIST_MAX = 60  # 最多持久化最近 N 条
_JOB_FIELDS = {f.name for f in dataclasses.fields(JobRecord)}

def persist_jobs(jobs) -> None:
    """把 jobs(最近 N 条)落盘供重启恢复；剥掉 retry_ctx 里的 api_key(不存明文)；全程吞异常。
    jobs: 调用方传入的 JobRecord 列表快照（调用方负责加锁）。"""
    try:
        jobs = list(jobs)[:QUEUE_PERSIST_MAX]
        out = []
        for j in jobs:
            d = dataclasses.asdict(j)
            ctx = d.get('retry_ctx')
            if isinstance(ctx, dict) and 'api_key' in ctx:
                ctx = dict(ctx); ctx.pop('api_key', None); d['retry_ctx'] = ctx
            if isinstance(ctx, dict):
                ctx = dict(ctx)
                for key in ('cpt', 'cpt_pro', 'sd_positive', 'sd_negative'):
                    value = ctx.pop(key, '')
                    if value:
                        ctx[f'{key}_obf'] = _obfuscate(value)
                d['retry_ctx'] = ctx
            out.append(d)
        tmp = QUEUE_STATE_FILE + '.tmp'
        with open(tmp, 'w', encoding='utf-8') as f:
            json.dump(out, f, ensure_ascii=False)
        os.replace(tmp, QUEUE_STATE_FILE)
    except Exception as ex:
        logger.warning(f"[队列] 持久化失败(忽略): {ex}")

def load_persisted_jobs() -> List[JobRecord]:
    """启动时读回队列；把中断态(queued/running)修正为 partial/failed。返回 JobRecord 列表。"""
    try:
        if not os.path.exists(QUEUE_STATE_FILE):
            return []
        with open(QUEUE_STATE_FILE, 'r', encoding='utf-8') as f:
            data = json.load(f)
        if not isinstance(data, list):
            return []
    except Exception as ex:
        logger.warning(f"[队列] 读取持久化失败(忽略): {ex}")
        return []
    jobs = []
    for d in data:
        if not isinstance(d, dict):
            continue
        try:
            job = JobRecord(**{k: v for k, v in d.items() if k in _JOB_FIELDS})
        except Exception:
            continue
        # 中断态修正：程序重启时仍标 queued/running 的任务不可能还在跑
        if job.status in ('queued', 'running'):
            generic_paths = any((r or {}).get('paths') for r in (job.model_runs or {}).values() if isinstance(r, dict))
            job.status = 'partial' if (job.b2_path or job.pro_path or generic_paths) else 'failed'
            job.error = '程序重启，任务已中断'
        job.b2_stage = ''; job.pro_stage = ''; job.pro_polishing = False
        for run in (job.model_runs or {}).values():
            if isinstance(run, dict):
                run['stage'] = ''
                if run.get('status') in ('queued', 'running'):
                    run['status'] = 'partial' if run.get('paths') else 'failed'
        if isinstance(job.retry_ctx, dict):
            for key in ('cpt', 'cpt_pro', 'sd_positive', 'sd_negative'):
                encoded = job.retry_ctx.pop(f'{key}_obf', '')
                if encoded and not job.retry_ctx.get(key):
                    job.retry_ctx[key] = _deobfuscate(encoded)
        ensure_candidate_lists(job)  # 老持久化只有 *_path → 回填成单元素候选列表(单张无 nav)，之后重抽即累积
        from .models import ensure_model_runs
        ensure_model_runs(job)
        jobs.append(job)
    if jobs:
        logger.info(f"[队列] 恢复 {len(jobs)} 个历史任务")
    return jobs


# ── 历史地板小样扫描（_ng_uploads 里的原始上传，供 webui 历史小样 picker）─────
_FLOOR_SWATCH_EXTS = ('.jpg', '.jpeg', '.png', '.webp')
# 历史小样里要排除的非地板上传：房间图/参照图/记录管理上传/测试临时图/匿名兜底名
_FLOOR_SWATCH_SKIP_PREFIX = ('room_', 'ref_', 'mgr_a_', 'mgr_b_', 'ZZ', 'upload_', 'logo_')

def _list_recent_floor_swatches(limit: int = 24):
    """扫 _ng_uploads 里的历史地板小样(原始上传)，按最近修改时间倒序返回绝对路径列表。
    文件名 stem 与 output_files/{材料}/ 文件夹同名 → 选用后自动复用同一材料文件夹、记录归并。"""
    items = []
    try:
        for f in os.listdir(UPLOAD_DIR):
            if not f.lower().endswith(_FLOOR_SWATCH_EXTS):
                continue
            if f.startswith(_FLOOR_SWATCH_SKIP_PREFIX):
                continue
            p = os.path.join(UPLOAD_DIR, f)
            if os.path.isfile(p):
                items.append((p, os.path.getmtime(p)))
    except Exception as ex:
        logger.warning(f"扫描历史小样失败: {ex}")
        return []
    items.sort(key=lambda x: x[1], reverse=True)
    return [p for p, _ in items[:limit]]


# ── 用量统计 ──────────────────────────────────────────────────────
# 累计「不同模式 × 模型(B2/Pro) × 线路(google/fal)」各出了多少张图、失败多少次。
# 不估算金额（计价多变不好控）。落盘 usage_stats.json，原子写，全程吞异常——绝不能拖垮生图主流程。
_USAGE_STATS_FILE = os.path.join(MAIN_OUTPUT_DIR, "usage_stats.json")
_usage_lock = threading.Lock()

def _short_mode_label(mode: str) -> str:
    """长 workflow_mode → 短标签（与记录页 split 显示一致）。"""
    return (mode or "").split("(")[0].split(" ")[0].strip() or "未知"

def _short_model_label(model: str) -> str:
    m = model or ""
    if "Aura" in m: return "AuraSR"
    if "SD35" in m or "SD 3.5" in m: return "SD35"
    if "Lite" in m: return "Lite"
    if "Pro" in m: return "Pro"
    if "B2" in m or " 2" in m or m.endswith("2"): return "B2"
    return m.strip() or "未知"

def _load_usage_raw() -> dict:
    def _normalize_preview_model(data: dict) -> dict:
        """旧版把 NB2 Lite 归并成 B2；按 operation=preview 可无歧义迁回 Lite。"""
        for operations in (data.get('counts') or {}).values():
            if not isinstance(operations, dict):
                continue
            preview = operations.get('preview')
            if not isinstance(preview, dict) or 'B2' not in preview:
                continue
            old = preview.get('B2')
            if not isinstance(old, dict):
                continue
            lite = preview.get('Lite')
            if not isinstance(lite, dict):
                lite = {}
                preview['Lite'] = lite
            preview.pop('B2', None)
            for provider, counts in old.items():
                if not isinstance(counts, dict):
                    continue
                row = lite.setdefault(provider, {'ok': 0, 'fail': 0})
                row['ok'] = int(row.get('ok', 0)) + int(counts.get('ok', 0))
                row['fail'] = int(row.get('fail', 0)) + int(counts.get('fail', 0))
        return data

    try:
        if os.path.exists(_USAGE_STATS_FILE):
            with open(_USAGE_STATS_FILE, "r", encoding="utf-8") as f:
                d = json.load(f)
            if isinstance(d, dict):
                if int(d.get('version', 1) or 1) >= 2:
                    return _normalize_preview_model(d)
                old = d.get('counts') or {}
                return {'version': 2, 'counts': {
                    mode: {'generate': models} for mode, models in old.items()
                    if isinstance(models, dict)
                }}
    except Exception as ex:
        logger.warning(f"[用量] 读取失败(将重置): {ex}")
    return {"version": 2, "counts": {}}

def _save_usage_raw(data: dict) -> None:
    tmp = _USAGE_STATS_FILE + ".tmp"
    with open(tmp, "w", encoding="utf-8") as f:
        json.dump(data, f, ensure_ascii=False)
    os.replace(tmp, _USAGE_STATS_FILE)

def record_usage(mode: str, model: str, provider: str, ok: bool, operation: str = 'generate') -> None:
    """累加一次出图结果。维度: 模式 × 模型 × 线路(google/fal/local)，各计 ok/fail。
    全程吞异常——统计绝不能影响生图。"""
    try:
        mkey = _short_mode_label(mode)
        mdl = _short_model_label(model)
        prov = (provider or "google").strip().lower()
        if prov == 'comfyui':
            prov = 'local'
        if prov not in ("google", "fal", "local"):
            prov = "unknown"
        with _usage_lock:
            data = _load_usage_raw()
            counts = data.setdefault("counts", {})
            op = (operation or 'generate').strip().lower()
            row = counts.setdefault(mkey, {}).setdefault(op, {}).setdefault(mdl, {}).setdefault(prov, {"ok": 0, "fail": 0})
            row["ok" if ok else "fail"] = int(row.get("ok" if ok else "fail", 0)) + 1
            _save_usage_raw(data)
    except Exception as ex:
        logger.debug(f"[用量] record_usage 忽略: {ex}")

def load_usage_summary(prices: Optional[dict] = None) -> dict:
    """结构化汇总，供 UI 渲染：
    {'rows': [{mode, model, provider, ok, fail, cost?}...(已排序)], 'totals': {'ok','fail','total','cost'}}
    prices: 可选单价表 {'B2': 0.1, 'Pro': 0.5, 'Lite': 0.03, 'B2:fal': 0.12,...}（元/张成功图）。
    行单价先查 '模型:线路' 再回落 '模型'；查无单价的行 cost=None（UI 显示 —），估算口径=只按成功计。"""
    data = _load_usage_raw()
    prices = prices or {}
    rows = []
    tot_ok = tot_fail = 0
    tot_cost = 0.0
    has_cost = False
    unpriced_ok = 0
    for mode, operations in sorted((data.get("counts") or {}).items()):
        if not isinstance(operations, dict):
            continue
        for operation, models in sorted(operations.items()):
            if not isinstance(models, dict):
                continue
            for model, provs in sorted(models.items()):
                if not isinstance(provs, dict):
                    continue
                for prov, c in sorted(provs.items()):
                    if not isinstance(c, dict):
                        continue
                    ok = int(c.get("ok", 0)); fail = int(c.get("fail", 0))
                    tot_ok += ok; tot_fail += fail
                    price = prices.get(f"{model}:{prov}", prices.get(model))
                    # 用量页统计 API 成本；本地 ComfyUI 未配置单价时按 0，而不是误报“未计价”。
                    if price is None and prov == 'local':
                        price = 0.0
                    cost = round(ok * price, 2) if price is not None else None
                    if cost is not None:
                        tot_cost += cost
                        has_cost = True
                    else:
                        unpriced_ok += ok
                    rows.append({"mode": mode, "operation": operation, "model": model,
                                 "provider": prov, "ok": ok, "fail": fail, "cost": cost})
    return {"rows": rows,
            "totals": {"ok": tot_ok, "fail": tot_fail, "total": tot_ok + tot_fail,
                       "cost": round(tot_cost, 2) if has_cost else None,
                       "unpriced_ok": unpriced_ok,
                       "cost_complete": unpriced_ok == 0}}

def _load_records(json_path):
    try:
        with open(json_path, 'r', encoding='utf-8') as f: return json.load(f)
    except json.JSONDecodeError as e:
        # 文件损坏：改名备份留待人工抢救。若原样留下，下一次保存会用空列表把整个历史覆盖掉。
        backup = f"{json_path}.corrupt_{time.strftime('%Y%m%d_%H%M%S')}"
        try:
            os.replace(json_path, backup)
            logger.error(f"记录文件损坏，已备份到 {backup} / {e}")
        except OSError as be:
            logger.error(f"记录文件损坏且备份失败: {json_path} / {e} / 备份错误: {be}")
        return []
    except Exception as e:
        if json_path and os.path.exists(json_path):
            logger.error(f"记录读取失败: {json_path} / {e}")
        return []

def room_type_counts(records) -> dict:
    """纯函数：把一组记录按非空 room_type 累计计数 → {房间类型: 张数}。空/缺失的 room_type 忽略。"""
    counts = {}
    for r in records or []:
        if not isinstance(r, dict):
            continue
        rt = (r.get('room_type', '') or '').strip()
        if rt:
            counts[rt] = counts.get(rt, 0) + 1
    return counts


def _save_records(json_path, records):
    # 先写同目录临时文件再原子替换：写一半崩溃/断电不会截断原文件
    tmp_path = None
    try:
        dir_path = os.path.dirname(json_path) or '.'
        fd, tmp_path = tempfile.mkstemp(prefix='.records_', suffix='.tmp', dir=dir_path)
        with os.fdopen(fd, 'w', encoding='utf-8') as f:
            json.dump(records, f, ensure_ascii=False, indent=2)
        # Windows: os.replace 撞上无锁读者(记录页浏览/导出/收藏扫描短暂 open 同文件)会
        # PermissionError；读者都是瞬时打开，写侧短重试即可，不给所有读者加锁。
        for _i in range(5):
            try:
                os.replace(tmp_path, json_path)
                break
            except PermissionError:
                if _i == 4:
                    raise
                time.sleep(0.1)
        tmp_path = None
    except Exception as e:
        logger.error(f"记录保存失败: {json_path} / {e}")
        raise
    finally:
        if tmp_path and os.path.exists(tmp_path):
            try: os.remove(tmp_path)
            except OSError: pass

def _delete_record(json_path, record_id):
    """删除 JSON 文件中指定 ID 的记录"""
    with record_file_lock(json_path):
        records = _load_records(json_path)
        new_records = [r for r in records if r.get('id') != record_id]
        if len(new_records) == len(records):
            logger.warning(f"[记录] 删除失败，未找到记录 json={json_path}, record={record_id}")
            return False
        _save_records(json_path, new_records)
    logger.info(f"[记录] 已删除记录 json={json_path}, record={record_id}")
    return True

def _delete_result_image(json_path, record_id, result_ref):
    """Delete a result reference by stable id (or a legacy integer index)."""
    with record_file_lock(json_path):
        records = _load_records(json_path)
        for r in records:
            if r.get('id') == record_id:
                results = r.get('results', [])
                idx = _result_index(results, result_ref)
                if idx >= 0:
                    results.pop(idx)
                    _save_records(json_path, records)
                    logger.info(f"[记录] 已删除效果图 json={json_path}, record={record_id}, result={result_ref}")
                    return True
    logger.warning(f"[记录] 删除效果图失败 json={json_path}, record={record_id}, result={result_ref}")
    return False


def _new_result_id() -> str:
    return f"res_{uuid.uuid4().hex}"


def _result_index(results, result_ref) -> int:
    """Resolve a stable result id; integers remain accepted for internal legacy callers."""
    if isinstance(result_ref, int):
        return result_ref if 0 <= result_ref < len(results) else -1
    return next((i for i, item in enumerate(results)
                 if item.get('result_id') == str(result_ref)), -1)


def _ensure_result_ids(records) -> bool:
    changed = False
    for record in records or []:
        for result in record.get('results', []) if isinstance(record, dict) else []:
            if not result.get('result_id'):
                result['result_id'] = _new_result_id()
                changed = True
    return changed

def _img_to_b64(img_or_path, max_width: Optional[int] = None) -> str:
    try:
        img = Image.open(img_or_path) if isinstance(img_or_path, str) else img_or_path.copy()
        if max_width and img.width > max_width: img = img.resize((max_width, int(img.height * max_width / img.width)), Image.Resampling.LANCZOS)
        if img.mode != 'RGB': img = img.convert('RGB')
        buf = io.BytesIO(); img.save(buf, format='JPEG', quality=85)
        return b64mod.b64encode(buf.getvalue()).decode('utf-8')
    except Exception as e:
        logger.error(f"图片转 base64 失败: {img_or_path} / {e}")
        return ''

def _b64_to_pil(b64_str):
    if not b64_str: return None
    try: return Image.open(io.BytesIO(b64mod.b64decode(b64_str)))
    except Exception as e:
        logger.error(f"base64 转图片失败: {e}")
        return None

def _rel_result_path(abs_path):
    """效果图绝对路径 → 相对 MAIN_OUTPUT_DIR 的路径(存进 result_image_file)。
    文件不存在或不在 output_files 下返回空串(空串时调用方回退 base64)。"""
    if not abs_path or not os.path.exists(abs_path):
        return ''
    try:
        rel = os.path.relpath(abs_path, MAIN_OUTPUT_DIR)
    except Exception:
        return ''
    return '' if rel.startswith('..') else rel


def _safe_output_path(rel):
    """把 result_image_file 相对路径解析为 MAIN_OUTPUT_DIR 内的绝对路径；越界/不存在/异常返回 None。
    正常路径由 _rel_result_path 生成(已挡 '..'),此处是对手改或历史污染 JSON 的二次兜底。"""
    if not rel:
        return None
    try:
        base = os.path.realpath(MAIN_OUTPUT_DIR)
        p = os.path.realpath(os.path.join(base, rel))
        # commonpath 在跨盘符(Windows 多盘)时抛 ValueError → 视为越界
        if os.path.commonpath([base, p]) != base:
            return None
    except (ValueError, OSError):
        return None
    return p if os.path.exists(p) else None


def scan_json_files():
    results = []
    if not os.path.exists(MAIN_OUTPUT_DIR): return results
    for root, dirs, files in os.walk(MAIN_OUTPUT_DIR):
        for f in files:
            if f.endswith('_记录.json'): results.append(os.path.join(root, f))
    return sorted(results, reverse=True)

def get_record_labels(json_path, records=None):
    """返回记录选项标签。records 可由调用方传入，避免文件列表汇总时重复读盘。"""
    if records is None:
        records = _load_records(json_path)
    choices = []
    for r in records:
        label = (f"{r.get('timestamp','')} | "
                 f"{r.get('workflow_mode','').split(' ')[0]} | "
                 f"{r.get('room_type','')}")
        choices.append((label, r.get('id', '')))
    return choices

def _html_b64(obj, file_key, b64_key):
    """HTML 导出取图：优先从文件读回 base64(保持 HTML 自包含)，回退内联 base64。"""
    rel = obj.get(file_key, '')
    p = _safe_output_path(rel)
    if p:
        try:
            with open(p, 'rb') as f:
                return b64mod.b64encode(f.read()).decode()
        except Exception as e:
            logger.warning(f"[导出] 读取图片失败 {rel}: {e}")
    return obj.get(b64_key, '')


def export_html_from_json(json_path):
    records = _load_records(json_path)
    if not records: return "❌ 没有找到记录"
    html_path = json_path.replace('_记录.json', '_导出.html')
    entries = []
    for r in records:
        _sb = _html_b64(r, 'sample_image_file', 'sample_image_b64')
        sample_tag = (f'<img src="data:image/jpeg;base64,{_sb}" '
                      f'style="width:180px;border-radius:6px;margin-bottom:10px;display:block;" />'
                      ) if _sb else ''
        results_html = ''
        for i, res in enumerate(r.get('results', [])):
            cmt = (f'<div style="margin-top:8px;padding:8px 12px;background:#fff8f0;border-left:3px solid #e8874a;border-radius:4px;">'
                   f'<p style="margin:0;font-size:0.82em;color:#7a4010;">💬 备注：{html.escape(res.get("comment",""))}</p></div>'
                   ) if res.get('comment') else ''
            review_bits = []
            if res.get('best'):
                review_bits.append('最佳图')
            _status = {'pass': '通过', 'backup': '备选', 'rejected': '淘汰', 'unreviewed': '未评'}.get(
                res.get('review_status', 'unreviewed'), '未评')
            if _status != '未评':
                review_bits.append(_status)
            if res.get('review_tags'):
                review_bits.append('标签：' + '、'.join(html.escape(str(t)) for t in res.get('review_tags', [])))
            if res.get('review_note'):
                review_bits.append('评审备注：' + html.escape(res.get('review_note', '')))
            review = (f'<div style="margin-top:8px;padding:8px 12px;background:#f3fbf8;border-left:3px solid #2e8c7e;border-radius:4px;">'
                      f'<p style="margin:0;font-size:0.82em;color:#235f55;">{" ｜ ".join(review_bits)}</p></div>'
                      ) if review_bits else ''
            _rb = _html_b64(res, 'result_image_file', 'result_image_b64')
            results_html += (f'<div style="margin-top:12px;padding:12px;background:#fdf8f4;border-radius:8px;">'
                             f'<p style="margin:0 0 8px 0;font-size:0.8em;color:#b07040;">📸 效果图 {i+1} — {res.get("result_timestamp","")}</p>'
                             f'<img src="data:image/jpeg;base64,{_rb}" style="max-width:100%;border-radius:6px;" />'
                             f'{cmt}{review}</div>')
        entries.append(f'<div style="padding:20px;border-bottom:3px solid #e8874a;margin-bottom:10px;">'
                       f'<p style="margin:0 0 8px 0;font-size:0.85em;color:#555;"><strong>🕒 {r.get("timestamp","")}</strong>'
                       f' &nbsp;|&nbsp; {r.get("workflow_mode","")} &nbsp;|&nbsp; {r.get("room_type","")}</p>'
                       f'{sample_tag}'
                       f'<pre style="white-space:pre-wrap;font-size:0.82em;background:#f8f8f8;padding:12px;border-radius:6px;">'
                       f'{html.escape(r.get("params_summary", ""))}</pre>'
                       f'{results_html}</div>')
    full_html = ('<!DOCTYPE html><html><head><meta charset="utf-8"><title>地板效果图记录</title>'
                 '<style>body{font-family:sans-serif;max-width:960px;margin:0 auto;padding:20px;}</style>'
                 f'</head><body>{"".join(entries)}</body></html>')
    with open(html_path, 'w', encoding='utf-8') as f: f.write(full_html)
    return f"✅ 已导出：{os.path.basename(html_path)}"

def append_result_to_log(img1_path, img2_path, json_path, record_id, comment1="", comment2=""):
    if not img1_path and not img2_path: return "⚠️ 请至少上传一张效果图"
    if not json_path or not record_id: return "⚠️ 请先加载一条记录"
    try:
        with record_file_lock(json_path):
            records = _load_records(json_path)
            for r in records:
                if r.get('id') == record_id:
                    written = []; ts = time.strftime("%Y-%m-%d %H:%M:%S")
                    for img_path, comment, label in [(img1_path, comment1, "Banana2"), (img2_path, comment2, "Pro")]:
                        if img_path:
                            rel = ''
                            try:
                                rel = _rel_result_path(_save_api_result_jpg(Image.open(img_path), label, json_path.replace('_记录.json', '_优化图.png')))
                            except Exception as ex:
                                logger.warning(f"[追加] 落盘失败(回退 base64): {ex}")
                            entry = {
                                'result_id': _new_result_id(),
                                'result_timestamp': ts,
                                'comment': str(comment).strip() if comment else '',
                                'model_label': label,
                            }
                            if rel:
                                entry['result_image_file'] = rel
                            else:
                                entry['result_image_b64'] = _img_to_b64(img_path, max_width=1000)
                            r.setdefault('results', []).append(entry)
                            written.append(label)
                    _save_records(json_path, records)
                    logger.info(
                        f"[记录] 手动追加效果图 json={json_path}, record={record_id}, written={written}, "
                        f"img1={img1_path}, img2={img2_path}"
                    )
                    return f"✅ 已写入：{' + '.join(written)}"
        logger.error(f"[记录] 手动追加失败，未找到记录 json={json_path}, record={record_id}")
        return "❌ 未找到对应记录"
    except Exception as e:
        logger.exception(f"[记录] 手动追加写入失败 json={json_path}, record={record_id}")
        return f"❌ 写入失败: {e}"

# ── 数据安全层 ────────────────────────────────────────────────
_PROMPT_KEY = b"braag2026floor_engine_v5_xor"

# 密码哈希加载优先级：环境变量 > 配置文件 > 内置默认值（仅开发用）
_DEFAULT_REVEAL_HASH = "455c459b728d459e5acf0373c929afc894ddb049515cd88cb046945e235e279e"

def _load_reveal_hash() -> str:
    """Load the reveal-password hash from env var, config, or built-in default."""
    import os as _os
    env_hash = _os.environ.get('FLOOR_ENGINE_REVEAL_HASH', '').strip()
    if env_hash:
        return env_hash
    cfg = _load_config()
    cfg_hash = cfg.get('reveal_hash', '').strip()
    if cfg_hash:
        return cfg_hash
    return _DEFAULT_REVEAL_HASH

def _obfuscate(text: str) -> str:
    if not text: return ""
    return b64mod.b64encode(bytes([b ^ _PROMPT_KEY[i % len(_PROMPT_KEY)] for i, b in enumerate(text.encode("utf-8"))])).decode()

def _deobfuscate(encoded: str) -> str:
    if not encoded: return ""
    try: return bytes([b ^ _PROMPT_KEY[i % len(_PROMPT_KEY)] for i, b in enumerate(b64mod.b64decode(encoded))]).decode("utf-8")
    except Exception: return ""


def migrate_record_storage(json_path) -> bool:
    """Remove plaintext prompts and add stable result ids, keeping one recovery backup."""
    with record_file_lock(json_path):
        records = _load_records(json_path)
        if not isinstance(records, list):
            return False
        changed = _ensure_result_ids(records)
        for record in records:
            if not isinstance(record, dict):
                continue
            prompt = record.pop('prompt_en', '')
            prompt_pro = record.pop('prompt_en_pro', '')
            if prompt:
                record.setdefault('_pe', _obfuscate(prompt))
                changed = True
            if prompt_pro:
                record.setdefault('_pe_pro', _obfuscate(prompt_pro))
                changed = True
            if record.get('_schema_version') != 2:
                record['_schema_version'] = 2
                changed = True
        if not changed:
            return False
        backup = json_path + '.schema_v1.bak'
        if not os.path.exists(backup) and os.path.exists(json_path):
            shutil.copy2(json_path, backup)
        _save_records(json_path, records)
        return True


def migrate_all_record_storage() -> int:
    changed = 0
    for path in scan_json_files():
        try:
            changed += int(migrate_record_storage(path))
        except Exception as ex:
            logger.error(f"[记录迁移] 失败，已保留原文件 {path}: {ex}")
    return changed

def reveal_prompt_fn(json_path, record_id, input_password):
    if not input_password: return "🔒 请输入密码"
    if hashlib.sha256(input_password.strip().encode('utf-8')).hexdigest() != _load_reveal_hash(): return "❌ 密码错误"
    for r in _load_records(json_path):
        if r.get('id') == record_id:
            pe = r.get('_pe', '')
            return _deobfuscate(pe) if pe else r.get('prompt_en', '无提示词数据')
    return "❌ 未找到记录"

# ── API 生成结果落地 & 记录写入（原在 api.py，归入记录层）──────────────
def _api_write_to_record(pil_img, model_key: str, json_path_val: str, record_id_val: str,
                         image_file: Optional[str] = None, metadata: Optional[dict] = None):
    # 用 is not None 而不是 bool(pil_img)——PIL Image 在某些版本布尔求值会异常
    if pil_img is None or not json_path_val or not record_id_val:
        logger.warning(f"_api_write_to_record 参数不全: img={pil_img is not None}, jpath={bool(json_path_val)}, rid={bool(record_id_val)}")
        return
    try:
        # 效果图文件化：优先存已落盘 jpg 的相对路径(记录页走 HTTP，不卡)；没传则自落盘兜底；
        # 实在不行才回退内联 base64(保证不丢图)。编码/落盘在锁外做，锁内只做读-改-写。
        rel = _rel_result_path(image_file)
        if not rel:
            rel = _rel_result_path(_save_api_result_jpg(pil_img, model_key, json_path_val.replace('_记录.json', '_优化图.png')))
        entry = {
            'result_id': _new_result_id(),
            'result_timestamp': time.strftime("%Y-%m-%d %H:%M:%S"),
            'comment': f'API 自动生成 ({pil_img.width}×{pil_img.height})',
            'model_label': model_key,
        }
        if metadata:
            entry['generation_metadata'] = dict(metadata)
        if rel:
            entry['result_image_file'] = rel
        else:
            entry['result_image_b64'] = _img_to_b64(pil_img, max_width=None)
        matched = False
        with record_file_lock(json_path_val):
            records = _load_records(json_path_val)
            for r in records:
                if r.get('id') == record_id_val:
                    r.setdefault('results', []).append(entry)
                    _save_records(json_path_val, records)
                    matched = True
                    break
        if matched:
            logger.info(f"✅ 已写入记录 {record_id_val} / {model_key} ({'file' if rel else 'b64'})")
            return entry['result_id']
        else:
            logger.warning(f"未找到记录 ID {record_id_val}，文件: {json_path_val}")
    except Exception as e:
        logger.error(f"❌ 写入记录失败 ({model_key}): {e}")
    return None

def _save_api_result_jpg(pil_img, model_key: str, png_path_val: str) -> str:
    tmp_path = None
    try:
        ts = time.strftime("%Y%m%d_%H%M%S") + f"_{int(time.time_ns() % 1_000_000_000):09d}"
        safe_key = ''.join(c if (c.isalnum() or c in '._-') else '_' for c in str(model_key or 'result')).strip('._') or 'result'
        raw_base = os.path.basename(png_path_val or "api")
        orig_name = raw_base.replace('_优化图.png', '').replace('.png', '') or "result"
        orig_name = ''.join(c if (c.isalnum() or c in '._-') else '_' for c in orig_name).strip('._') or 'result'
        fname = f"{orig_name}_{safe_key}_{ts}_{uuid.uuid4().hex[:10]}.jpg"
        fpath = os.path.join(MAIN_OUTPUT_DIR, fname)
        img = pil_img.copy()
        if img.mode != 'RGB': img = img.convert('RGB')
        fd, tmp_path = tempfile.mkstemp(prefix='.result_', suffix='.jpg', dir=MAIN_OUTPUT_DIR)
        os.close(fd)
        img.save(tmp_path, format='JPEG', quality=95)
        os.replace(tmp_path, fpath)
        tmp_path = None
        return fpath
    except Exception as e:
        logger.error(f"API 结果图片保存失败 ({model_key}): {e}")
        return None
    finally:
        if tmp_path and os.path.exists(tmp_path):
            try: os.remove(tmp_path)
            except OSError: pass


def _save_api_result_png(pil_img, model_key: str, source_path: str,
                         metadata: Optional[dict] = None) -> str:
    """给 SD/超分结果保存无损 PNG；命名与 JPG 出口同样防碰撞。"""
    tmp_path = None
    try:
        ts = time.strftime("%Y%m%d_%H%M%S") + f"_{int(time.time_ns() % 1_000_000_000):09d}"
        safe_key = ''.join(c if (c.isalnum() or c in '._-') else '_' for c in str(model_key or 'result')).strip('._') or 'result'
        raw_base = os.path.basename(source_path or "api")
        orig_name = raw_base.replace('_优化图.png', '').rsplit('.', 1)[0] or "result"
        orig_name = ''.join(c if (c.isalnum() or c in '._-') else '_' for c in orig_name).strip('._') or 'result'
        fpath = os.path.join(MAIN_OUTPUT_DIR, f"{orig_name}_{safe_key}_{ts}_{uuid.uuid4().hex[:10]}.png")
        fd, tmp_path = tempfile.mkstemp(prefix='.result_', suffix='.png', dir=MAIN_OUTPUT_DIR)
        os.close(fd)
        pnginfo = None
        if metadata:
            pnginfo = PngImagePlugin.PngInfo()
            pnginfo.add_text('floor_engine', json.dumps(metadata, ensure_ascii=False, separators=(',', ':')))
        pil_img.convert('RGB').save(tmp_path, format='PNG', optimize=True, pnginfo=pnginfo)
        os.replace(tmp_path, fpath)
        tmp_path = None
        return fpath
    except Exception as ex:
        logger.error(f"API PNG 结果保存失败 ({model_key}): {ex}")
        return None
    finally:
        if tmp_path and os.path.exists(tmp_path):
            try: os.remove(tmp_path)
            except OSError: pass

def append_edited_result_to_record(json_path, record_id, source_ref, pil_img, edit_prompt, model_label, image_file=None):
    if pil_img is None:
        return "❌ 没有可写入的二次修改图片"
    if not json_path or not record_id:
        return "⚠️ 请先加载一条记录"
    try:
        rel = _rel_result_path(image_file)
        if not rel:
            rel = _rel_result_path(_save_api_result_jpg(pil_img, f'{model_label}_Edit', json_path.replace('_记录.json', '_优化图.png')))
        with record_file_lock(json_path):
            records = _load_records(json_path)
            for r in records:
                if r.get('id') == record_id:
                    ts = time.strftime("%Y-%m-%d %H:%M:%S")
                    comment = (
                        f"二次修改自结果 {source_ref} | {model_label}\n"
                        f"修改建议：{str(edit_prompt).strip()}"
                    )
                    entry = {
                        'result_id': _new_result_id(),
                        'result_timestamp': ts,
                        'comment': comment,
                        'model_label': f"{model_label} Edit",
                        'source_result_id': str(source_ref),
                        'edit_prompt': str(edit_prompt).strip(),
                    }
                    if rel:
                        entry['result_image_file'] = rel
                    else:
                        entry['result_image_b64'] = _img_to_b64(pil_img, max_width=None)
                    r.setdefault('results', []).append(entry)
                    _save_records(json_path, records)
                    logger.info(
                        f"[二改记录] 已追加 record={record_id}, source={source_ref}, "
                        f"model={model_label}, image={pil_img.width}x{pil_img.height}, instruction_len={len(str(edit_prompt or ''))}"
                    )
                    return "✅ 二次修改结果已追加到当前记录"
        logger.error(f"[二改记录] 未找到对应记录 json={json_path}, record={record_id}")
        return "❌ 未找到对应记录"
    except Exception as e:
        logger.exception(f"[二改记录] 写入失败 json={json_path}, record={record_id}")
        return f"❌ 写入失败: {e}"


# ── 收藏（给满意的效果图打星）────────────────────────────────────
def toggle_result_favorite(json_path, record_id, result_ref):
    """翻转某条结果的收藏标记，返回新状态(True/False)；未找到返回 None。"""
    with record_file_lock(json_path):
        records = _load_records(json_path)
        for r in records:
            if r.get('id') == record_id:
                results = r.get('results', [])
                idx = _result_index(results, result_ref)
                if idx >= 0:
                    new_state = not bool(results[idx].get('favorite'))
                    results[idx]['favorite'] = new_state
                    _save_records(json_path, records)
                    logger.info(f"[记录] 收藏切换 json={json_path}, record={record_id}, result={result_ref}, fav={new_state}")
                    return new_state
    logger.warning(f"[记录] 收藏切换失败，未找到 json={json_path}, record={record_id}, result={result_ref}")
    return None


def attach_generation_context(json_path, record_id, ctx: dict) -> bool:
    """把一次生成的完整入参快照挂到记录上（gen_context），供「复用参数」「前后对比」使用。
    ctx 由调用方保证不含密钥。找不到记录返回 False。"""
    with record_file_lock(json_path):
        records = _load_records(json_path)
        for r in records:
            if r.get('id') == record_id:
                r['gen_context'] = ctx
                _save_records(json_path, records)
                return True
    logger.warning(f"[记录] 生成上下文写入失败，未找到 json={json_path}, record={record_id}")
    return False


REVIEW_STATUSES = {'unreviewed', 'pass', 'backup', 'rejected'}


def update_result_review(json_path, record_id, result_ref, *,
                         review_status='unreviewed', review_tags=None,
                         review_note='', best=False):
    """写入单张效果图的人工评审元数据。

    best=True 时，同一 record 下其它结果的 best 会被清掉，保证“最佳图”唯一。
    返回更新后的评审快照；找不到记录/索引返回 None。
    """
    status = (review_status or 'unreviewed').strip()
    if status not in REVIEW_STATUSES:
        status = 'unreviewed'
    tags = []
    for t in review_tags or []:
        s = str(t or '').strip()
        if s and s not in tags:
            tags.append(s)
    note = str(review_note or '').strip()
    best = bool(best)

    with record_file_lock(json_path):
        records = _load_records(json_path)
        for r in records:
            if r.get('id') != record_id:
                continue
            results = r.get('results', [])
            idx = _result_index(results, result_ref)
            if idx < 0:
                break
            if best:
                for res in results:
                    res['best'] = False
            target = results[idx]
            target['review_status'] = status
            target['review_tags'] = tags
            target['review_note'] = note
            target['best'] = best
            target['reviewed_at'] = time.strftime("%Y-%m-%d %H:%M:%S")
            _save_records(json_path, records)
            logger.info(
                f"[记录] 评审更新 json={json_path}, record={record_id}, "
                f"result={result_ref}, status={status}, tags={tags}, best={best}"
            )
            return {
                'review_status': status,
                'review_tags': tags,
                'review_note': note,
                'best': best,
                'reviewed_at': target['reviewed_at'],
            }
    logger.warning(f"[记录] 评审更新失败 json={json_path}, record={record_id}, result={result_ref}")
    return None


def collect_favorites():
    """扫描所有记录文件，汇总被收藏(favorite=True)的结果，按出图时间倒序。
    返回 [{json_path, material, record_id, record, result_id, res}, ...]。"""
    out = []
    for jp in scan_json_files():
        material = os.path.basename(jp).replace('_记录.json', '').replace('.json', '')
        try:
            recs = _load_records(jp)
        except Exception:
            continue
        for r in recs:
            for res in r.get('results', []):
                if res.get('favorite'):
                    out.append({
                        'json_path': jp, 'material': material,
                        'record_id': r.get('id', ''), 'record': r,
                        'result_id': res.get('result_id', ''), 'res': res,
                    })
    out.sort(key=lambda x: x['res'].get('result_timestamp', ''), reverse=True)
    return out


# ── 评审复盘：人工评审标签的聚合统计与好图样本库 ─────────────────────
_REVIEW_DIM_FIELDS = ('workflow_mode', 'style', 'room_type', 'seam')


def load_review_summary() -> dict:
    """扫描全部记录，按维度聚合人工评审结果，供复盘面板渲染。
    通过率口径 = pass / 已评审数(pass+backup+rejected)，未评审不摊薄；行无已评审时 pass_rate=None。
    返回 {'overview': {...}, 'dimensions': {dim: [行...]}, 'tags': [{'tag','count'}...]}。"""
    dims = {d: {} for d in _REVIEW_DIM_FIELDS}
    tag_counter = {}
    total = passed = backup = rejected = best_count = 0
    for jp in scan_json_files():
        for r in _load_records(jp):
            if not isinstance(r, dict):
                continue
            keys = {
                'workflow_mode': _short_mode_label(r.get('workflow_mode', '')),
                'style': (str(r.get('style') or '')).strip() or '未知',
                'room_type': (str(r.get('room_type') or '')).strip() or '未知',
                'seam': (str(r.get('seam') or '')).strip() or '未知',
            }
            for res in r.get('results', []):
                if not isinstance(res, dict):
                    continue
                status = (res.get('review_status') or 'unreviewed')
                if status not in REVIEW_STATUSES:
                    status = 'unreviewed'
                total += 1
                if status == 'pass':
                    passed += 1
                elif status == 'backup':
                    backup += 1
                elif status == 'rejected':
                    rejected += 1
                if res.get('best'):
                    best_count += 1
                for t in res.get('review_tags') or []:
                    t = str(t or '').strip()
                    if t:
                        tag_counter[t] = tag_counter.get(t, 0) + 1
                for dim, key in keys.items():
                    row = dims[dim].setdefault(
                        key, {'pass': 0, 'backup': 0, 'rejected': 0, 'unreviewed': 0, 'total': 0})
                    row[status] += 1
                    row['total'] += 1

    def _rate(p, b, rj):
        done = p + b + rj
        return round(p / done, 3) if done else None

    dimensions = {}
    for dim, rows in dims.items():
        out_rows = []
        for key, c in rows.items():
            out_rows.append({'key': key, **c,
                             'pass_rate': _rate(c['pass'], c['backup'], c['rejected'])})
        out_rows.sort(key=lambda x: x['total'], reverse=True)
        dimensions[dim] = out_rows

    reviewed = passed + backup + rejected
    return {
        'overview': {
            'total': total, 'reviewed': reviewed,
            'coverage': round(reviewed / total, 3) if total else None,
            'pass': passed, 'backup': backup, 'rejected': rejected,
            'pass_rate': _rate(passed, backup, rejected),
            'best': best_count,
        },
        'dimensions': dimensions,
        'tags': sorted(({'tag': t, 'count': n} for t, n in tag_counter.items()),
                       key=lambda x: x['count'], reverse=True),
    }


def collect_review_gallery(filter_key: str = 'pass', limit: int = 60):
    """好图样本库：filter_key='pass'(评审通过) 或 'best'(最佳图)，按出图时间倒序取 limit 条。
    返回 [{json_path, material, record_id, result_id, res, style, room_type, workflow_mode}, ...]。"""
    out = []
    for jp in scan_json_files():
        material = os.path.basename(jp).replace('_记录.json', '').replace('.json', '')
        for r in _load_records(jp):
            if not isinstance(r, dict):
                continue
            for res in r.get('results', []):
                if not isinstance(res, dict):
                    continue
                hit = res.get('best') if filter_key == 'best' else (
                    res.get('review_status') == 'pass')
                if hit:
                    out.append({
                        'json_path': jp, 'material': material,
                        'record_id': r.get('id', ''), 'result_id': res.get('result_id', ''),
                        'res': res,
                        'style': r.get('style', ''), 'room_type': r.get('room_type', ''),
                        'workflow_mode': _short_mode_label(r.get('workflow_mode', '')),
                    })
    out.sort(key=lambda x: x['res'].get('result_timestamp', ''), reverse=True)
    return out[:max(1, int(limit or 60))]


# ── PPTX 导出（客户提案 deck）──────────────────────────────────────
def _result_image_source(res):
    """结果图供 PPTX 用的来源：优先成图文件绝对路径，回退 base64 解成 BytesIO；无图返回 None。"""
    p = _safe_output_path(res.get('result_image_file', ''))
    if p:
        return p
    b64 = res.get('result_image_b64', '')
    if b64:
        try:
            return io.BytesIO(b64mod.b64decode(b64))
        except Exception as e:
            logger.warning(f"[PPTX] base64 解码失败: {e}")
    return None


def _build_pptx(items, out_path, title, branding=None):
    """items: [{'source': 路径或BytesIO, 'caption': 文本}]。16:9 PPTX：标题页 + 每项一页(图+说明)。
    python-pptx 只存字体名(微软雅黑)，PowerPoint 用系统字体渲染中文，无需嵌字体。
    branding: 可选 {'company','contact','logo_path'}（config.get_pptx_branding()）——
    标题页加 logo/公司名/联系方式，内容页右下角加公司名页脚；配置坏/图打不开一律静默降级。"""
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt, Emu
        from pptx.enum.text import PP_ALIGN
        from pptx.dml.color import RGBColor
    except Exception:
        return "❌ 未安装 python-pptx，请先运行: pip install python-pptx"
    from PIL import Image as _PILImage

    branding = branding or {}
    company = str(branding.get('company') or '').strip()
    contact = str(branding.get('contact') or '').strip()
    logo_path = str(branding.get('logo_path') or '').strip()

    prs = Presentation()
    prs.slide_width = Inches(13.333); prs.slide_height = Inches(7.5)  # 16:9
    SW, SH = int(prs.slide_width), int(prs.slide_height)
    blank = prs.slide_layouts[6]

    def _set_text(text_frame, text, size, bold, color, align=None):
        p = text_frame.paragraphs[0]
        if align is not None: p.alignment = align
        run = p.add_run(); run.text = text
        run.font.size = Pt(size); run.font.bold = bold
        run.font.name = '微软雅黑'; run.font.color.rgb = color

    # 标题页（logo 置顶居中 → 主标题 → 公司名副标题 → 联系方式底部小字）
    s0 = prs.slides.add_slide(blank)
    if logo_path and os.path.isfile(logo_path):
        try:
            with _PILImage.open(logo_path) as im:
                lw, lh = im.size
            # Pillow 给的是像素，python-pptx 使用 EMU。这里的 scale 单位是
            # EMU/像素，不能再限制为 <= 1，否则普通 1000px logo 只会得到
            # 1000 EMU（约 0.001 英寸），在封面上几乎不可见。
            max_w, max_h = int(Inches(2.4)), int(Inches(1.2))
            scale = min(max_w / max(1, lw), max_h / max(1, lh))
            pw, ph = int(lw * scale), int(lh * scale)
            s0.shapes.add_picture(logo_path, int((SW - pw) / 2), int(Inches(1.2)),
                                  width=pw, height=ph)
        except Exception as e:
            logger.warning(f"[PPTX] logo 插入失败(忽略): {e}")
    tb = s0.shapes.add_textbox(Inches(0.8), Inches(2.8), SW - Inches(1.6), Inches(1.8))
    tb.text_frame.word_wrap = True
    _set_text(tb.text_frame, title, 40, True, RGBColor(0x2b, 0x24, 0x1a), PP_ALIGN.CENTER)
    if company:
        sb = s0.shapes.add_textbox(Inches(0.8), Inches(4.4), SW - Inches(1.6), Inches(0.7))
        sb.text_frame.word_wrap = True
        _set_text(sb.text_frame, company, 18, False, RGBColor(0x6b, 0x63, 0x56), PP_ALIGN.CENTER)
    if contact:
        cb = s0.shapes.add_textbox(Inches(0.8), SH - Inches(1.0), SW - Inches(1.6), Inches(0.6))
        cb.text_frame.word_wrap = True
        _set_text(cb.text_frame, contact, 12, False, RGBColor(0x9a, 0x90, 0x82), PP_ALIGN.CENTER)

    area_w = SW - int(Inches(1.0)); area_h = SH - int(Inches(2.0))
    top0 = int(Inches(0.4))
    for it in items:
        src = it.get('source')
        if src is None:
            continue
        slide = prs.slides.add_slide(blank)
        try:
            if hasattr(src, 'seek'): src.seek(0)
            with _PILImage.open(src) as im: iw, ih = im.size
            if hasattr(src, 'seek'): src.seek(0)
        except Exception:
            iw, ih = 4, 3
        scale = min(area_w / iw, area_h / ih)
        pw = int(iw * scale); ph = int(ih * scale)
        left = int((SW - pw) / 2)
        try:
            slide.shapes.add_picture(src, left, top0, width=Emu(pw), height=Emu(ph))
        except Exception as e:
            logger.warning(f"[PPTX] 插图失败(跳过一页): {e}")
            continue
        cap = slide.shapes.add_textbox(Inches(0.5), SH - Inches(1.4), SW - Inches(1.0), Inches(1.2))
        cap.text_frame.word_wrap = True
        _set_text(cap.text_frame, it.get('caption', ''), 12, False, RGBColor(0x55, 0x55, 0x55))
        if company:
            ft = slide.shapes.add_textbox(SW - Inches(3.2), SH - Inches(0.42),
                                          Inches(3.0), Inches(0.35))
            _set_text(ft.text_frame, company, 9, False, RGBColor(0xb0, 0xa8, 0x9a), PP_ALIGN.RIGHT)

    try:
        prs.save(out_path)
    except Exception as e:
        logger.exception(f"[PPTX] 保存失败 {out_path}")
        return f"❌ 保存失败: {e}"
    logger.info(f"[PPTX] 已导出 {out_path}（{len(items)} 张）")
    return f"✅ 已导出：{os.path.basename(out_path)}"


def _result_caption(material, res, record):
    caption = " · ".join(p for p in [
        material, res.get('model_label', ''), res.get('result_timestamp', ''),
        (record.get('params_summary', '') or ''),
    ] if p)
    cmt = res.get('comment', '')
    review = []
    status = {'pass': '通过', 'backup': '备选', 'rejected': '淘汰', 'unreviewed': ''}.get(
        res.get('review_status', 'unreviewed'), '')
    if res.get('best'):
        review.append('最佳图')
    if status:
        review.append(status)
    if res.get('review_tags'):
        review.append('标签：' + '、'.join(str(t) for t in res.get('review_tags', [])))
    if res.get('review_note'):
        review.append('评审备注：' + str(res.get('review_note', '')))
    tail = (f"\n备注：{cmt}" if cmt else "")
    if review:
        tail += "\n评审：" + " ｜ ".join(review)
    return caption + tail


def export_pptx_from_json(json_path):
    """把当前材料的所有效果图导出成一份 PPTX。"""
    records = _load_records(json_path)
    if not records:
        return "❌ 没有找到记录"
    material = os.path.basename(json_path).replace('_记录.json', '').replace('.json', '')
    items = []
    for r in records:
        for res in r.get('results', []):
            src = _result_image_source(res)
            if src is not None:
                items.append({'source': src, 'caption': _result_caption(material, res, r)})
    if not items:
        return "❌ 该文件没有可导出的效果图"
    out_path = json_path.replace('_记录.json', '_导出.pptx')
    return _build_pptx(items, out_path, f"{material} · 地板效果图", branding=get_pptx_branding())


def export_favorites_pptx():
    """把所有收藏(跨材料)的效果图合成一份客户提案 PPTX。"""
    favs = collect_favorites()
    if not favs:
        return "❌ 还没有收藏任何效果图（先在记录里点 ⭐）"
    items = []
    for f in favs:
        src = _result_image_source(f['res'])
        if src is not None:
            items.append({'source': src, 'caption': _result_caption(f['material'], f['res'], f['record'])})
    if not items:
        return "❌ 收藏的效果图都没有可用图片数据"
    out_path = os.path.join(MAIN_OUTPUT_DIR, f"收藏夹提案_{time.strftime('%Y%m%d_%H%M%S')}.pptx")
    return _build_pptx(items, out_path, "地板效果图 · 收藏提案", branding=get_pptx_branding())


def migrate_record_file(json_path) -> bool:
    """把记录里的 base64 图片导出成文件、改存 *_image_file、删 base64，给 JSON 瘦身、记录页提速。
    迁移前备份 .bak；幂等(已迁移则返回 False)；单张失败则保留其 base64 容错。"""
    def _need(recs):
        for r in recs:
            if r.get('sample_image_b64') and not r.get('sample_image_file'):
                return True
            for res in r.get('results', []):
                if res.get('result_image_b64') and not res.get('result_image_file'):
                    return True
        return False
    try:
        if not _need(_load_records(json_path)):
            return False
    except Exception:
        return False
    with record_file_lock(json_path):
        records = _load_records(json_path)
        if not _need(records):
            return False
        bak = json_path + '.bak'
        if not os.path.exists(bak):
            try: shutil.copy2(json_path, bak)
            except Exception as e: logger.warning(f"[迁移] 备份失败(继续): {e}")
        pseudo_png = os.path.basename(json_path).replace('_记录.json', '_优化图.png')
        changed = 0
        for r in records:
            sb = r.get('sample_image_b64')
            if sb and not r.get('sample_image_file'):
                pil = _b64_to_pil(sb)
                if pil is not None:
                    rel = _rel_result_path(_save_api_result_jpg(pil, 'sample', pseudo_png))
                    if rel:
                        r['sample_image_file'] = rel
                        r.pop('sample_image_b64', None); changed += 1
            for res in r.get('results', []):
                b = res.get('result_image_b64')
                if b and not res.get('result_image_file'):
                    pil = _b64_to_pil(b)
                    if pil is not None:
                        rel = _rel_result_path(_save_api_result_jpg(pil, res.get('model_label') or 'result', pseudo_png))
                        if rel:
                            res['result_image_file'] = rel
                            res.pop('result_image_b64', None); changed += 1
        if changed:
            _save_records(json_path, records)
        logger.info(f"[迁移] {os.path.basename(json_path)} 完成，迁移 {changed} 张图，备份={os.path.basename(bak)}")
        return changed > 0


__all__ = [
    'record_file_lock', '_get_json_path', '_load_records', '_save_records',
    '_delete_record', '_delete_result_image', '_img_to_b64', '_b64_to_pil',
    'scan_json_files', 'get_record_labels', 'export_html_from_json',
    'append_result_to_log', 'append_edited_result_to_record',
    '_obfuscate', '_deobfuscate', 'reveal_prompt_fn',
    '_api_write_to_record', '_save_api_result_jpg', 'migrate_record_file',
    'migrate_record_storage', 'migrate_all_record_storage',
    'toggle_result_favorite', 'update_result_review', 'collect_favorites',
    'attach_generation_context', 'load_review_summary', 'collect_review_gallery',
    'export_pptx_from_json', 'export_favorites_pptx',
]
