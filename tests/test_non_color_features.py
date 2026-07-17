import json
import math
import os

import pytest
from fastapi import HTTPException
from PIL import Image
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

from Floor_engine_server import server_schemas, server_helpers, routes_config, routes_library
from Floor_engine_server import config, records, server_api, usage_stats, exports


def test_pptx_brand_logo_uses_real_slide_units(tmp_path):
    logo = tmp_path / "logo.png"
    output = tmp_path / "branded.pptx"
    Image.new("RGB", (1000, 400), (180, 40, 30)).save(logo)

    msg = exports._build_pptx([], str(output), "测试", {"logo_path": str(logo)})

    assert msg.startswith("✅")
    prs = Presentation(output)
    pictures = [
        shape for shape in prs.slides[0].shapes
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE
    ]
    assert len(pictures) == 1
    picture = pictures[0]
    assert picture.width / 914400 == pytest.approx(2.4, abs=0.01)
    assert picture.height / 914400 == pytest.approx(0.96, abs=0.01)


def test_usage_prices_lite_separately_and_marks_partial_total(tmp_path, monkeypatch):
    usage_file = tmp_path / "usage.json"
    usage_file.write_text(json.dumps({
        "version": 2,
        "counts": {
            "纯效果图": {
                # 历史文件里 Lite 曾被错误归并为 B2；读取时应无损迁回 Lite。
                "preview": {"B2": {"google": {"ok": 5, "fail": 0}}},
                "generate": {
                    "B2": {"google": {"ok": 2, "fail": 0}},
                    "Pro": {"google": {"ok": 1, "fail": 0}},
                },
            },
        },
    }), encoding="utf-8")
    monkeypatch.setattr(usage_stats, "_USAGE_STATS_FILE", str(usage_file))

    summary = usage_stats.load_usage_summary({"B2": 1.0, "Lite": 0.1})
    rows = {(row["operation"], row["model"]): row for row in summary["rows"]}

    assert usage_stats._short_model_label("NB2 Lite") == "Lite"
    assert rows[("preview", "Lite")]["cost"] == 0.5
    assert rows[("generate", "B2")]["cost"] == 2.0
    assert rows[("generate", "Pro")]["cost"] is None
    assert summary["totals"]["cost"] == 2.5
    assert summary["totals"]["unpriced_ok"] == 1
    assert summary["totals"]["cost_complete"] is False


def test_local_inpaint_usage_is_not_charged_as_google(tmp_path, monkeypatch):
    usage_file = tmp_path / "usage.json"
    monkeypatch.setattr(usage_stats, "_USAGE_STATS_FILE", str(usage_file))

    usage_stats.record_usage("纯效果图", "ComfyUI", "comfyui", True, "inpaint")
    summary = usage_stats.load_usage_summary({})

    assert summary["rows"] == [{
        "mode": "纯效果图", "operation": "inpaint", "model": "ComfyUI",
        "provider": "local", "ok": 1, "fail": 0, "cost": 0.0,
    }]
    assert summary["totals"]["cost"] == 0.0
    assert summary["totals"]["unpriced_ok"] == 0


def test_usage_prices_reject_non_finite_values(tmp_path, monkeypatch):
    cfg_file = tmp_path / "engine_config.json"
    cfg_file.write_text('{"usage_prices":{"B2":Infinity,"Pro":NaN,"Lite":0.2}}', encoding="utf-8")
    monkeypatch.setattr(config, "CONFIG_FILE", str(cfg_file))

    assert config.get_usage_prices() == {"Lite": 0.2}
    with pytest.raises(HTTPException) as exc:
        routes_config.put_config(server_schemas.ConfigPatch(usage_prices={"B2": math.inf}))
    assert exc.value.status_code == 400


def test_clear_logo_removes_only_managed_upload(tmp_path, monkeypatch):
    uploads = tmp_path / "uploads"
    uploads.mkdir()
    logo = uploads / "logo_old.png"
    logo.write_bytes(b"old")
    patches = []
    monkeypatch.setattr(server_helpers, "UPLOAD_DIR", str(uploads))
    monkeypatch.setattr(
        routes_library,
        "get_pptx_branding",
        lambda: {"company": "", "contact": "", "logo_path": str(logo)},
    )
    monkeypatch.setattr(routes_library, "update_config", lambda patch: patches.append(patch) or True)

    assert routes_library.clear_logo() == {"ok": True}
    assert patches == [{"pptx_logo_path": ""}]
    assert not os.path.exists(logo)
